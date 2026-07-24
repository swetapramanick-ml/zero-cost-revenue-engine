import os
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
IMAGES_DIR = ROOT / "static" / "images"
IMAGES_DIR.mkdir(parents=True, exist_ok=True)


def create_technical_images():
    """Generate simple technical diagrams for the presentation."""
    pipeline_path = IMAGES_DIR / "pipeline_flow.png"
    architecture_path = IMAGES_DIR / "system_architecture.png"

    if not pipeline_path.exists():
        image = Image.new("RGB", (1200, 700), (248, 250, 252))
        draw = ImageDraw.Draw(image)
        draw.rectangle((40, 40, 1160, 660), outline=(31, 78, 121), width=4)
        boxes = [
            (120, 220, "Ingest Lead"),
            (360, 220, "Extract\nMetadata"),
            (620, 220, "AI\nPersonalize"),
            (900, 220, "Approve\n& Send"),
        ]
        for x, y, text in boxes:
            draw.rounded_rectangle((x, y, x + 180, y + 140), radius=18, outline=(46, 109, 164), width=3)
            draw.text((x + 90, y + 64), text, fill=(31, 78, 121), anchor="mm")
        draw.line((300, 290, 360, 290), fill=(46, 109, 164), width=4)
        draw.line((540, 290, 620, 290), fill=(46, 109, 164), width=4)
        draw.line((800, 290, 900, 290), fill=(46, 109, 164), width=4)
        draw.polygon([(300, 290), (315, 275), (315, 305)], fill=(46, 109, 164))
        draw.polygon([(540, 290), (555, 275), (555, 305)], fill=(46, 109, 164))
        draw.polygon([(800, 290), (815, 275), (815, 305)], fill=(46, 109, 164))
        image.save(pipeline_path)

    if not architecture_path.exists():
        image = Image.new("RGB", (1200, 700), (248, 250, 252))
        draw = ImageDraw.Draw(image)
        draw.rectangle((80, 100, 1120, 600), outline=(31, 78, 121), width=4)
        boxes = [
            (140, 170, "Frontend\nHTML / CSS / JS"),
            (420, 170, "FastAPI\nREST API"),
            (700, 170, "SQLite\nDatabase"),
            (980, 170, "Gemini + SMTP"),
        ]
        for x, y, text in boxes:
            draw.rounded_rectangle((x, y, x + 180, y + 140), radius=18, outline=(46, 109, 164), width=3)
            draw.text((x + 90, y + 64), text, fill=(31, 78, 121), anchor="mm")
        draw.line((320, 240, 420, 240), fill=(46, 109, 164), width=4)
        draw.line((600, 240, 700, 240), fill=(46, 109, 164), width=4)
        draw.line((880, 240, 980, 240), fill=(46, 109, 164), width=4)
        draw.polygon([(320, 240), (335, 225), (335, 255)], fill=(46, 109, 164))
        draw.polygon([(600, 240), (615, 225), (615, 255)], fill=(46, 109, 164))
        draw.polygon([(880, 240), (895, 225), (895, 255)], fill=(46, 109, 164))
        draw.text((240, 430), "Lead ingestion → enrichment → AI draft → approval → dispatch", fill=(31, 78, 121))
        image.save(architecture_path)

    return pipeline_path, architecture_path


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 78, 121)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(2.2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(230, 240, 250)
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, content_list, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    title_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(31, 78, 121)
    title_shape.line.color.rgb = RGBColor(31, 78, 121)

    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_before = Pt(8)

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.6))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, content in enumerate(content_list):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(19)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.level = 0

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.1), Inches(8.6), Inches(0.7))
        sub_frame = sub_box.text_frame
        p = sub_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(80, 100, 120)
        p.alignment = PP_ALIGN.RIGHT

    return slide


def add_code_slide(prs, title, code_lines, image_path=None, image_caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    title_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(31, 78, 121)
    title_shape.line.color.rgb = RGBColor(31, 78, 121)

    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_before = Pt(8)

    code_bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.1), Inches(5.7), Inches(5.8))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = RGBColor(244, 247, 250)
    code_bg.line.color.rgb = RGBColor(31, 78, 121)

    code_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(5.3), Inches(5.4))
    code_frame = code_box.text_frame
    code_frame.word_wrap = True
    for i, line in enumerate(code_lines):
        p = code_frame.paragraphs[0] if i == 0 else code_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.name = "Consolas"
        p.font.color.rgb = RGBColor(15, 23, 42)
        p.space_after = Pt(2)

    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(str(image_path), Inches(6.5), Inches(1.45), width=Inches(3.0), height=Inches(2.9))
        if image_caption:
            caption_box = slide.shapes.add_textbox(Inches(6.45), Inches(4.45), Inches(3.1), Inches(0.7))
            caption_frame = caption_box.text_frame
            p = caption_frame.paragraphs[0]
            p.text = image_caption
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(80, 100, 120)
            p.alignment = PP_ALIGN.CENTER

    return slide


def create_presentation(filename="Zero_Cost_Revenue_Engine_Presentation.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    pipeline_image, architecture_image = create_technical_images()

    add_title_slide(
        prs,
        "Zero-Cost Revenue Engine",
        "AI-Powered Sales Outreach Automation\nSweta Pramanick | 25500123176 | Dr Sudhir Chandra Sur Institute of Technology & Sports Complex"
    )

    add_content_slide(
        prs,
        "Project Overview",
        [
            "Objective: automate B2B outreach with AI-driven personalization and zero-cost infrastructure.",
            "Pipeline: lead ingestion → metadata extraction → AI drafting → human review → SMTP dispatch.",
            "Value: saves hundreds of manual hours while keeping quality control in the loop.",
            "Demo: https://zero-cost-revenue-engine.onrender.com/"
        ],
        subtitle="Built with FastAPI, Gemini API, SQLite, and a lightweight dashboard"
    )

    add_content_slide(
        prs,
        "Why This Project Matters",
        [
            "Traditional outreach is slow, repetitive, and expensive for startups and freelancers.",
            "Generic templates often underperform because they lack contextual relevance.",
            "This solution combines web scraping, structured storage, LLM-based personalization, and approval workflows.",
            "The result is a practical, low-cost growth engine for real sales teams."
        ]
    )

    add_code_slide(
        prs,
        "Core Pipeline Logic",
        [
            "lead = db.get_lead(lead_id)",
            "metadata = extract_domain_metadata(lead['domain'])",
            "db.update_lead_enrichment(lead_id, metadata, status='Enriched')",
            "ai_data = generate_personalized_content(company, domain, metadata)",
            "db.update_lead_ai_output(lead_id, ai_data, status='Pending_Review')"
        ],
        image_path=str(pipeline_image),
        image_caption="Lead processing workflow"
    )

    add_code_slide(
        prs,
        "AI Personalization Example",
        [
            "prompt = f'''You are a top-tier B2B sales rep...'''",
            "response = client.models.generate_content(",
            "    model='gemini-3.5-flash',",
            "    contents=prompt,",
            "    config=types.GenerateContentConfig(...)",
            ")",
            "return json.loads(response.text)"
        ],
        image_path=str(architecture_image),
        image_caption="System architecture overview"
    )

    add_content_slide(
        prs,
        "Key Components",
        [
            "FastAPI backend with async endpoints for ingesting leads, updating drafts, approving sends, and managing settings.",
            "SQLite database stores lead state, extracted metadata, AI-generated copy, and SMTP configuration.",
            "The frontend dashboard provides a human-in-the-loop handoff for review, edit, approve, or reject actions."
        ]
    )

    add_content_slide(
        prs,
        "Database Design",
        [
            "Leads table: id, company_name, domain, industry, extracted_metadata, pain_points, recipient_email, personalized_subject, personalized_body, status.",
            "Settings table: GEMINI_API_KEY, SMTP_HOST, SMTP_PORT, SMTP_USER, SMTP_PASSWORD, SMTP_FROM_EMAIL, SMTP_FROM_NAME.",
            "This design makes the pipeline stateful, auditable, and easy to extend."
        ]
    )

    add_code_slide(
        prs,
        "SMTP Dispatch Snippet",
        [
            "if port == 465:",
            "    server = smtplib.SMTP_SSL(host, port)",
            "else:",
            "    server = smtplib.SMTP(host, port)",
            "    server.starttls()",
            "server.login(user, password)",
            "server.sendmail(from_email, [to_email], msg.as_string())"
        ]
    )

    add_content_slide(
        prs,
        "Technical Highlights",
        [
            "Free-tier AI integration keeps the project affordable while still producing useful personalization.",
            "Background-task processing lets the system handle multiple leads without blocking the UI.",
            "The modular design allows swapping the scraper, model, or transporter without rewriting the workflow."
        ]
    )

    add_content_slide(
        prs,
        "Live Demo and Repository",
        [
            "Demo URL: https://zero-cost-revenue-engine.onrender.com/",
            "GitHub: https://github.com/swetapramanick-ml/zero-cost-revenue-engine",
            "The project includes the full FastAPI app, database layer, extractor, AI logic, and SMTP dispatcher."
        ]
    )

    add_title_slide(prs, "Thank You!", "Questions and discussion")

    output_path = ROOT / filename
    prs.save(output_path)
    print("Presentation created successfully: " + str(output_path))


if __name__ == "__main__":
    create_presentation()
