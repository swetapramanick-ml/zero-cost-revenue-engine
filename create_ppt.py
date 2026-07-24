import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation(filename="Zero_Cost_Revenue_Engine_Presentation.pptx"):
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    # -----------------------------------------
    # Slide 1: Title Slide
    # -----------------------------------------
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Zero-Cost Revenue Engine"
    
    subtitle_text = (
        "My Name: Sweta Pramanick\n"
        "University Roll Number: 25500123176\n"
        "College Name: Dr Sudhir Chandra Sur Institute of Technology & Sports Complex\n"
        "Department: BTech in Computer Science & Engineering\n"
        "Session: 2023-27\n\n"
        "Live Demo: https://zero-cost-revenue-engine.onrender.com/\n"
        "GitHub: https://github.com/swetapramanick-ml/zero-cost-revenue-engine"
    )
    subtitle.text = subtitle_text
    
    # Format subtitle font size to fit well
    for p in subtitle.text_frame.paragraphs:
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER
        
    # -----------------------------------------
    # Helper function for body slides
    # -----------------------------------------
    def add_bullet_slide(prs, title_text, bullet_points):
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = title_text
        
        tf = body_shape.text_frame
        tf.text = bullet_points[0]
        
        for point in bullet_points[1:]:
            p = tf.add_paragraph()
            p.text = point
            if point.startswith("  -") or point.startswith("\t"):
                p.level = 1
                p.text = point.strip(" \t-")
            elif point.startswith("    -"):
                p.level = 2
                p.text = point.strip(" \t-")
                
        # Adjust font sizes
        for p in tf.paragraphs:
            if p.level == 0:
                p.font.size = Pt(24)
            elif p.level == 1:
                p.font.size = Pt(20)
            elif p.level == 2:
                p.font.size = Pt(18)
        
        return slide

    # -----------------------------------------
    # Slide 2: Introduction
    # -----------------------------------------
    add_bullet_slide(
        prs,
        "Introduction",
        [
            "What is the Zero-Cost Revenue Engine?",
            "  - An automated Sales Outreach Pipeline.",
            "  - Designed to eliminate manual sales prospecting.",
            "  - Extracts data from target domains and enriches it automatically.",
            "  - Leverages AI to generate highly personalized outreach emails.",
            "  - Maintains a 'Human-in-the-Loop' for quality assurance."
        ]
    )

    # -----------------------------------------
    # Slide 3: The Problem Statement
    # -----------------------------------------
    add_bullet_slide(
        prs,
        "The Problem Statement",
        [
            "Current sales outreach is highly inefficient:",
            "  - Manual Prospecting: Takes hours to research a single lead.",
            "  - Generic Emails: Low conversion rates due to lack of personalization.",
            "  - Expensive Tools: Existing sales CRMs and AI tools cost thousands of dollars.",
            "  - Human Error: Manual tracking of lead status leads to missed opportunities.",
            "Business Need: A streamlined, automated, and free tool to drive revenue."
        ]
    )

    # -----------------------------------------
    # Slide 4: Our Solution
    # -----------------------------------------
    add_bullet_slide(
        prs,
        "Our Solution: Zero-Cost Revenue Engine",
        [
            "A seamless automation pipeline built from open-source and free-tier tools.",
            "Fully automated lead ingestion and enrichment.",
            "AI-powered contextual personalization tailored to each lead's domain.",
            "Integrated Email Dispatcher to send drafts straight to clients.",
            "Achieves 'Zero-Cost' by utilizing free APIs (like Gemini) and lightweight frameworks."
        ]
    )

    # -----------------------------------------
    # Slide 5: System Architecture
    # -----------------------------------------
    add_bullet_slide(
        prs,
        "System Architecture",
        [
            "Frontend:",
            "  - Responsive Dashboard (HTML, CSS, Vanilla JS).",
            "Backend:",
            "  - FastAPI (Python) for robust, async API endpoints.",
            "Database:",
            "  - SQLite for maintaining lead state transitions.",
            "AI & Integrations:",
            "  - Google Gemini API for Natural Language Processing.",
            "  - SMTP protocol for automated email dispatch."
        ]
    )

    # -----------------------------------------
    # Slide 6: The 4-Phase Pipeline
    # -----------------------------------------
    add_bullet_slide(
        prs,
        "The 4-Phase Pipeline",
        [
            "The system operates in a background task queue with four main phases:",
            "  1. Data Extraction & Enrichment",
            "  2. AI Personalization",
            "  3. Human Gate Approval",
            "  4. Automated Dispatch",
            "Lead Status is tracked across: Ingested -> Enriched -> Pending Review -> Completed."
        ]
    )

    # -----------------------------------------
    # Slide 7: Phase 1 - Data Extraction
    # -----------------------------------------
    add_bullet_slide(
        prs,
        "Phase 1: Data Extraction",
        [
            "Goal: Gather context to write a highly targeted email.",
            "Process:",
            "  - Input: Lead's Company Name & Website Domain.",
            "  - Scraper automatically visits the domain.",
            "  - Extracts vital metadata: Meta descriptions, industry keywords, and contact info.",
            "  - Fallback mechanisms for identifying primary email addresses.",
            "Output: Enriched Lead Profile stored in the database."
        ]
    )

    # -----------------------------------------
    # Slide 8: Phase 2 - AI Personalization
    # -----------------------------------------
    add_bullet_slide(
        prs,
        "Phase 2: AI Personalization",
        [
            "Goal: Write an email that actually gets read.",
            "Integration: Google Gemini API (1.5 Flash/2.5 Flash).",
            "Process:",
            "  - Analyzes the enriched metadata.",
            "  - Identifies potential 'Pain Points' for the business.",
            "  - Generates a compelling, personalized subject line.",
            "  - Drafts a tailored email body offering a specific value proposition.",
            "Status updates to 'Pending Review'."
        ]
    )

    # -----------------------------------------
    # Slide 9: Phase 3 - Human Gate Approval
    # -----------------------------------------
    add_bullet_slide(
        prs,
        "Phase 3: Human Gate Approval",
        [
            "While automation is powerful, human intuition is essential.",
            "The Dashboard provides an interface to review AI-generated emails.",
            "Actions Available:",
            "  - Approve: Email looks great, ready to send.",
            "  - Edit: Tweak the subject line or body for perfection.",
            "  - Regenerate: Prompt the AI to try a different angle.",
            "  - Reject: Drop unqualified leads from the pipeline."
        ]
    )

    # -----------------------------------------
    # Slide 10: Phase 4 - Automated Dispatch
    # -----------------------------------------
    add_bullet_slide(
        prs,
        "Phase 4: Automated Dispatch",
        [
            "Once approved, the pipeline handles the final delivery.",
            "Features:",
            "  - Configurable SMTP settings (Host, Port, User, Password).",
            "  - Secure email transmission directly from the platform.",
            "  - Transaction-safe updates: Status changes to 'Completed' only if email sends successfully.",
            "  - Error logging and failure state management."
        ]
    )

    # -----------------------------------------
    # Slide 11: Technology Stack Highlights
    # -----------------------------------------
    add_bullet_slide(
        prs,
        "Technology Stack Highlights",
        [
            "FastAPI: Chosen for its speed, automatic interactive docs (Swagger UI), and background task support.",
            "Gemini API: Powerful, context-aware LLM that offers a generous free tier, maintaining the 'Zero-Cost' ethos.",
            "SQLite: Zero-configuration SQL database, perfect for portability and simple state management.",
            "Vanilla JS/HTML/CSS: Lightweight frontend without the overhead of heavy frameworks like React."
        ]
    )

    # -----------------------------------------
    # Slide 12: Core Advantages
    # -----------------------------------------
    add_bullet_slide(
        prs,
        "Core Advantages",
        [
            "Cost-Effective: Replaces expensive SaaS tools like Apollo or Lemlist.",
            "Time-Saving: Turns hours of prospecting into seconds of API calls.",
            "Highly Scalable: Background workers allow concurrent processing of multiple leads.",
            "Better Conversion Rates: Deeply personalized emails significantly outperform generic templates.",
            "Data Privacy: Your leads and AI prompts remain within your local pipeline."
        ]
    )

    # -----------------------------------------
    # Slide 13: Live Demonstration
    # -----------------------------------------
    add_bullet_slide(
        prs,
        "Live Demonstration",
        [
            "You can try the engine out right now:",
            "Demo URL: https://zero-cost-revenue-engine.onrender.com/",
            "Workflow to test:",
            "  1. Add a target company and domain.",
            "  2. Watch the status change from 'Ingested' to 'Pending Review'.",
            "  3. Review the AI-generated pain points and email copy.",
            "  4. Explore the Settings page for SMTP/API configurations."
        ]
    )

    # -----------------------------------------
    # Slide 14: Future Enhancements
    # -----------------------------------------
    add_bullet_slide(
        prs,
        "Future Enhancements",
        [
            "CRM Integration: Syncing directly with HubSpot or Salesforce.",
            "Multi-Channel Outreach: Expanding beyond email to LinkedIn automated messaging.",
            "Advanced Analytics: Tracking email open rates, click-through rates, and replies.",
            "Bulk Import: Support for uploading CSV lists of hundreds of domains.",
            "A/B Testing: Automatically testing different AI prompts to optimize conversion."
        ]
    )

    # -----------------------------------------
    # Slide 15: Conclusion & Q&A
    # -----------------------------------------
    add_bullet_slide(
        prs,
        "Conclusion & Q&A",
        [
            "Summary: The Zero-Cost Revenue Engine democratizes enterprise-grade sales automation.",
            "Thank You for your time and attention!",
            "",
            "Questions?",
            "",
            "Contact & Links:",
            "  - Demo: https://zero-cost-revenue-engine.onrender.com/",
            "  - GitHub: https://github.com/swetapramanick-ml/zero-cost-revenue-engine"
        ]
    )

    prs.save(filename)
    print(f"Presentation saved successfully as '{filename}'")

if __name__ == '__main__':
    create_presentation()
